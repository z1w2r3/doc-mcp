#!/usr/bin/env python
"""
分析样本 PPT 文件,了解其结构
"""
from pptx import Presentation
from pathlib import Path

# 样本文件路径
sample_file = Path("output/江苏省智能建造试点项目.pptx")

if not sample_file.exists():
    print(f"❌ 文件不存在: {sample_file}")
    exit(1)

print(f"📄 分析文件: {sample_file.name}")
print(f"📏 文件大小: {sample_file.stat().st_size / 1024:.1f} KB\n")

# 打开演示文稿
prs = Presentation(str(sample_file))

print(f"📊 基本信息:")
print(f"- 幻灯片总数: {len(prs.slides)}")
print(f"- 幻灯片尺寸: {prs.slide_width} x {prs.slide_height}")

# 提取核心属性
if prs.core_properties:
    print(f"\n📋 文档元数据:")
    props = prs.core_properties
    print(f"- 标题: {props.title or '(无)'}")
    print(f"- 作者: {props.author or '(无)'}")
    print(f"- 主题: {props.subject or '(无)'}")
    print(f"- 创建时间: {props.created or '(无)'}")
    print(f"- 修改时间: {props.modified or '(无)'}")
    print(f"- 最后修改人: {props.last_modified_by or '(无)'}")

# 分析每一张幻灯片
print(f"\n🎯 幻灯片详情:")
print("=" * 80)

for idx, slide in enumerate(prs.slides, 1):
    print(f"\n📄 幻灯片 #{idx}")
    print(f"- 形状数量: {len(slide.shapes)}")

    # 提取标题
    if slide.shapes.title:
        print(f"- 标题: {slide.shapes.title.text}")

    # 分析所有形状
    text_shapes = []
    table_shapes = []
    picture_shapes = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and shape != slide.shapes.title:  # 跳过标题(已显示)
                text_shapes.append(text[:50] + "..." if len(text) > 50 else text)

        if shape.has_table:
            table = shape.table
            table_shapes.append(f"表格 {len(table.rows)}行 x {len(table.columns)}列")

        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            picture_shapes.append("图片")

    if text_shapes:
        print(f"- 文本形状 ({len(text_shapes)}):")
        for i, text in enumerate(text_shapes[:3], 1):  # 只显示前3个
            print(f"  {i}. {text}")
        if len(text_shapes) > 3:
            print(f"  ... 还有 {len(text_shapes) - 3} 个")

    if table_shapes:
        print(f"- 表格 ({len(table_shapes)}): {', '.join(table_shapes)}")

    if picture_shapes:
        print(f"- 图片 ({len(picture_shapes)})")

    # 提取备注
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            print(f"- 备注: {notes_text[:100]}...")

print("\n" + "=" * 80)
print(f"✅ 分析完成!")
